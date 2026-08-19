# -*- coding: utf-8 -*-
"""Math Course AI — tanıtım sunumunu TEK KAYNAKTAN üretir.

İçerik `tanitim_icerik.py`'de; bu betik ondan sekiz dosya üretir:

    Math_Course_AI_Tanitim.pptx        · .pdf · .html      (TR, ekran)
    Math_Course_AI_Tanitim_Baski.pptx  · .pdf              (TR, baskı)
    Math_Course_AI_Intro_EN.pptx       · .pdf · .html      (EN, ekran)
    Math_Course_AI_Intro_EN_Print.pptx · .pdf              (EN, baskı)

TASARIM KARARI — neden şablon kopyalama yok
-------------------------------------------
Elle tasarlanmış bir kaynak sunum yoktu, bu yüzden slaytlar sıfırdan
kuruluyor ama tasarım TEK YERDE (`_Tasarim`): renk/ölçü/hizalama tahmin
edilmiyor,
programın kendi paletinden (mca_common.THEME / DARK_THEME) türetiliyor.
Böylece sunum ile uygulama aynı kimliği paylaşır ve tema değişirse sunum da
onunla birlikte değişir.

Yeniden çalıştırılabilirlik: her koşu BOŞ bir sunumdan başlar (python-pptx'te
slayt silmek pakette artık bırakabiliyor), dolayısıyla idempotency ispat
edilecek bir şey değil, TANIM GEREĞİ sağlanır.

PDF ve slayt PNG'leri PowerPoint COM ile üretilir; PowerPoint yoksa PPTX ve
HTML yine üretilir (HTML o durumda metin tabanlı slaytlarla çalışır) ve
eksik olan açıkça raporlanır — sessizce yarım çıktı verilmez.

Kullanım:
    python tanitim_uret.py                 # hepsini üret
    python tanitim_uret.py --sadece-pptx   # PDF/PNG/HTML yok
    python tanitim_uret.py --kontrol       # üretmeden durum raporu
"""

import base64
import glob
import html as _html
import os
import shutil
import sys
import tempfile

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from mca_common import VERSION
from tanitim_icerik import SLAYTLAR, ALTBASLIK

# ══════════════════════════════════════════════════════════════════════════
# YAZI TIPI — bilerek LIBRE bir aile secilir
# ══════════════════════════════════════════════════════════════════════════
# Uretilen PDF yeniden dagitilabilir olmali. Tescilli bir sistem yazi tipinin
# (ornegin Segoe UI) gomulmesi, o yazi tipinin kendi izin bitlerine baglidir;
# libre bir aile bu bagimliligi tamamen ortadan kaldirir. Sira: Inter →
# Source Sans 3 → DejaVu Sans. Hicbiri kurulu degilse geri dusulur ve bu
# DURUM ACIKCA RAPOR EDILIR (sessizce tescilli yazi tipine kacilmaz).
LIBRE_YAZI_TIPLERI = [
    ("Inter", "SIL Open Font License 1.1"),
    ("Source Sans 3", "SIL Open Font License 1.1"),
    ("Source Sans Pro", "SIL Open Font License 1.1"),
    ("DejaVu Sans", "DejaVu Fonts License (Bitstream Vera derivative)"),
]
YEDEK_YAZI_TIPI = ("Segoe UI", "Microsoft system font - proprietary, NOT libre")


def _kurulu_yazi_tipleri():
    """Windows'ta kurulu yazi tipi ailelerinin adlari (kucuk harf kumesi)."""
    adlar = set()
    try:
        import glob as _glob
        from pathlib import Path as _P
        kokler = [_P(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"]
        yerel = os.environ.get("LOCALAPPDATA")
        if yerel:
            kokler.append(_P(yerel) / "Microsoft" / "Windows" / "Fonts")
        for kok in kokler:
            if not kok.is_dir():
                continue
            for f in kok.iterdir():
                if f.suffix.lower() in (".ttf", ".otf", ".ttc"):
                    adlar.add(f.stem.lower().replace("-", " ").replace("_", " "))
    except Exception:
        pass
    try:                      # matplotlib bilir; kuruluysa daha guvenilir
        from matplotlib import font_manager as _fm
        for f in _fm.fontManager.ttflist:
            adlar.add(str(f.name).lower())
    except Exception:
        pass
    return adlar


def sec_yazi_tipi():
    """(aile_adi, lisans, libre_mi) — THIRD_PARTY_NOTICES.md bunu kaydeder."""
    kurulu = _kurulu_yazi_tipleri()
    for ad, lisans in LIBRE_YAZI_TIPLERI:
        anahtar = ad.lower()
        if any(anahtar in k for k in kurulu):
            return ad, lisans, True
    return YEDEK_YAZI_TIPI[0], YEDEK_YAZI_TIPI[1], False


YAZI_TIPI, YAZI_TIPI_LISANS, YAZI_TIPI_LIBRE = sec_yazi_tipi()

KOK = os.path.dirname(os.path.abspath(__file__))
GORSEL_KLASOR = os.path.join(KOK, "tanitim_gorseller")

CIKTILAR = {
    ("tr", "ekran"): "Math_Course_AI_Tanitim",
    ("tr", "baski"): "Math_Course_AI_Tanitim_Baski",
    ("en", "ekran"): "Math_Course_AI_Intro_EN",
    ("en", "baski"): "Math_Course_AI_Intro_EN_Print",
}

# ══════════════════════════════════════════════════════════════════════════
# TASARIM — tek yer. Ekran sürümü programın koyu teması, baskı sürümü açık
# teması. İkisi de mca_common paletinden türetilir ki sunum ile uygulama
# aynı kimliği paylaşsın.
# ══════════════════════════════════════════════════════════════════════════
class _Tasarim:
    def __init__(self, mod):
        from mca_common import THEME, DARK_THEME
        self.mod = mod
        if mod == "baski":
            p = dict(THEME)
            self.zemin = "FFFFFF"           # baskıda kâğıt beyazı: mürekkep yemesin
            self.panel = _hex(p["panel2"])
            self.metin = _hex(p["text"])
            self.soluk = "6B6154"
            self.cizgi = _hex(p["border"])
            self.vurgu = "B4551F"           # amber baskıda koyulaştırılır
            self.vurgu_metin = "FFFFFF"
            self.yesil = "0F7A42"
            self.kapak_zemin = "F6EFE4"
            self.kapak_metin = _hex(p["text"])
        else:
            p = dict(DARK_THEME)
            self.zemin = _hex(p["bg"])
            self.panel = _hex(p["panel"])
            self.metin = _hex(p["text"])
            self.soluk = _hex(p["muted"])
            self.cizgi = _hex(p["border"])
            self.vurgu = _hex(p["accent"])
            self.vurgu_metin = "FFFFFF"
            self.yesil = _hex(p["green"])
            self.kapak_zemin = _hex(p["surface"])
            self.kapak_metin = _hex(p["text"])


def _hex(deger):
    return str(deger).lstrip("#").upper()


def _rgb(hex6):
    return RGBColor.from_string(_hex(hex6))


GENISLIK = Inches(13.333)
YUKSEKLIK = Inches(7.5)
KENAR = Inches(0.62)


# ─────────────────────────────────────────────────────────── çizim yardımı ──
def _dikdortgen(slayt, x, y, w, h, dolgu, cizgi=None, yaricap=True):
    sekil = slayt.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if yaricap else MSO_SHAPE.RECTANGLE, x, y, w, h)
    sekil.fill.solid()
    sekil.fill.fore_color.rgb = _rgb(dolgu)
    if cizgi:
        sekil.line.color.rgb = _rgb(cizgi)
        sekil.line.width = Pt(1)
    else:
        sekil.line.fill.background()
    sekil.shadow.inherit = False
    if yaricap:
        try:                       # köşe yarıçapı: kart görünümü için yumuşak
            sekil.adjustments[0] = 0.045
        except Exception:
            pass
    return sekil


def _metin(slayt, x, y, w, h, parcalar, hiza=PP_ALIGN.LEFT, dikey=MSO_ANCHOR.TOP,
           satir_araligi=1.0):
    """parcalar: [(metin, punto, renk, kalin)] — her biri ayrı paragraf."""
    kutu = slayt.shapes.add_textbox(x, y, w, h)
    tf = kutu.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = dikey
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (yazi, punto, renk, kalin) in enumerate(parcalar):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = hiza
        p.line_spacing = satir_araligi
        for j, satir in enumerate(str(yazi).split("\n")):
            r = p.add_run()
            r.text = satir
            r.font.size = Pt(punto)
            r.font.bold = kalin
            r.font.color.rgb = _rgb(renk)
            r.font.name = YAZI_TIPI
            if j < len(str(yazi).split("\n")) - 1:
                p.add_line_break()
    return kutu


def _slayt(prs, t):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dikdortgen(s, 0, 0, GENISLIK, YUKSEKLIK, t.zemin, yaricap=False)
    return s


def _baslik_bandi(s, t, baslik, alt=None):
    _dikdortgen(s, 0, 0, GENISLIK, Inches(0.09), t.vurgu, yaricap=False)
    y = Inches(0.5)
    _metin(s, KENAR, y, GENISLIK - 2 * KENAR, Inches(0.75),
           [(baslik, 30, t.metin, True)])
    if alt:
        _metin(s, KENAR, Inches(1.22), GENISLIK - 2 * KENAR, Inches(0.6),
               [(alt, 14, t.soluk, False)])
    return Inches(2.0) if alt else Inches(1.58)


def _dipnot(s, t, yazi):
    if not yazi:
        return
    _metin(s, KENAR, YUKSEKLIK - Inches(0.72), GENISLIK - 2 * KENAR, Inches(0.45),
           [(yazi, 11, t.soluk, False)])


def _sayfa_no(s, t, no, toplam):
    _metin(s, GENISLIK - Inches(1.5), YUKSEKLIK - Inches(0.52), Inches(1.0),
           Inches(0.3), [(f"{no}/{toplam}", 10, t.soluk, False)],
           hiza=PP_ALIGN.RIGHT)


# ────────────────────────────────────────────────────────────── slayt tipleri ──
def _ciz_kapak(s, t, veri, dil):
    _dikdortgen(s, 0, 0, GENISLIK, YUKSEKLIK, t.kapak_zemin, yaricap=False)
    _dikdortgen(s, 0, 0, Inches(0.22), YUKSEKLIK, t.vurgu, yaricap=False)
    _dikdortgen(s, Inches(1.0), Inches(1.75), Inches(0.9), Inches(0.9), t.vurgu)
    _metin(s, Inches(1.0), Inches(1.92), Inches(0.9), Inches(0.6),
           [("Σ", 34, t.vurgu_metin, True)], hiza=PP_ALIGN.CENTER)
    _metin(s, Inches(2.2), Inches(1.75), Inches(10.0), Inches(1.2),
           [(_c(veri["baslik"], dil), 54, t.kapak_metin, True)])
    _metin(s, Inches(2.2), Inches(3.0), Inches(9.6), Inches(1.5),
           [(_c(veri["alt"], dil), 19, t.soluk, False)], satir_araligi=1.25)
    _dikdortgen(s, Inches(2.2), Inches(4.45), Inches(2.6), Inches(0.05), t.vurgu,
                yaricap=False)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_madde(s, t, veri, dil):
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    ogeler = veri["ogeler"]
    yukseklik = Inches(0.95)
    bosluk = Inches(0.2)
    for i, oge in enumerate(ogeler):
        y = ust + i * (yukseklik + bosluk)
        _dikdortgen(s, KENAR, y, Inches(0.06), yukseklik, t.vurgu, yaricap=False)
        _metin(s, KENAR + Inches(0.28), y, GENISLIK - 2 * KENAR - Inches(0.4),
               yukseklik, [(_c(oge, dil), 16, t.metin, False)],
               dikey=MSO_ANCHOR.MIDDLE, satir_araligi=1.18)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_kart(s, t, veri, dil):
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    ogeler = veri["ogeler"]
    sutun = 2
    w = (GENISLIK - 2 * KENAR - Inches(0.3)) / sutun
    h = Inches(1.95)
    for i, (bas, govde) in enumerate(ogeler):
        satir, sut = divmod(i, sutun)
        x = KENAR + sut * (w + Inches(0.3))
        y = ust + satir * (h + Inches(0.28))
        _dikdortgen(s, x, y, w, h, t.panel, t.cizgi)
        _dikdortgen(s, x, y, Inches(0.055), h, t.vurgu, yaricap=False)
        _metin(s, x + Inches(0.32), y + Inches(0.22), w - Inches(0.6), Inches(0.5),
               [(_c(bas, dil), 17, t.vurgu, True)])
        _metin(s, x + Inches(0.32), y + Inches(0.72), w - Inches(0.6),
               h - Inches(0.9), [(_c(govde, dil), 13, t.metin, False)],
               satir_araligi=1.2)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_sayi(s, t, veri, dil):
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    ogeler = veri["ogeler"]
    n = len(ogeler)
    w = (GENISLIK - 2 * KENAR - Inches(0.3) * (n - 1)) / n
    h = Inches(2.9)
    for i, (sayi, aciklama) in enumerate(ogeler):
        x = KENAR + i * (w + Inches(0.3))
        y = ust + Inches(0.4)
        _dikdortgen(s, x, y, w, h, t.panel, t.cizgi)
        _metin(s, x, y + Inches(0.55), w, Inches(1.1),
               [(_c(sayi, dil), 46, t.vurgu, True)], hiza=PP_ALIGN.CENTER)
        _metin(s, x + Inches(0.2), y + Inches(1.75), w - Inches(0.4), Inches(1.0),
               [(_c(aciklama, dil), 13, t.soluk, False)], hiza=PP_ALIGN.CENTER,
               satir_araligi=1.2)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_tablo(s, t, veri, dil):
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    satirlar = veri["ogeler"]
    h = Inches(0.62)
    w1 = Inches(5.2)
    w2 = GENISLIK - 2 * KENAR - w1 - Inches(0.2)
    for i, (sol, sag) in enumerate(satirlar):
        y = ust + i * (h + Inches(0.12))
        basliksa = i == 0
        _dikdortgen(s, KENAR, y, w1, h, t.vurgu if basliksa else t.panel,
                    None if basliksa else t.cizgi)
        _dikdortgen(s, KENAR + w1 + Inches(0.2), y, w2, h,
                    t.vurgu if basliksa else t.panel,
                    None if basliksa else t.cizgi)
        renk = t.vurgu_metin if basliksa else t.metin
        _metin(s, KENAR + Inches(0.25), y, w1 - Inches(0.4), h,
               [(_c(sol, dil), 15, renk, basliksa)], dikey=MSO_ANCHOR.MIDDLE)
        _metin(s, KENAR + w1 + Inches(0.45), y, w2 - Inches(0.4), h,
               [(_c(sag, dil), 15, renk, basliksa)], dikey=MSO_ANCHOR.MIDDLE)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_kapanis(s, t, veri, dil):
    _dikdortgen(s, 0, 0, GENISLIK, YUKSEKLIK, t.kapak_zemin, yaricap=False)
    _dikdortgen(s, 0, YUKSEKLIK - Inches(0.14), GENISLIK, Inches(0.14), t.vurgu,
                yaricap=False)
    _metin(s, KENAR, Inches(0.85), GENISLIK - 2 * KENAR, Inches(1.0),
           [(_c(veri["baslik"], dil), 40, t.kapak_metin, True)])
    for i, oge in enumerate(veri["ogeler"]):
        y = Inches(2.3) + i * Inches(1.15)
        _dikdortgen(s, KENAR, y, Inches(0.5), Inches(0.5), t.vurgu)
        _metin(s, KENAR, y + Inches(0.09), Inches(0.5), Inches(0.4),
               [(str(i + 1), 20, t.vurgu_metin, True)], hiza=PP_ALIGN.CENTER)
        _metin(s, KENAR + Inches(0.85), y, GENISLIK - 2 * KENAR - Inches(1.0),
               Inches(0.9), [(_c(oge, dil), 17, t.kapak_metin, False)],
               satir_araligi=1.2)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _gorsel_yolu(ad, dil):
    """tanitim_gorseller/<dil>/<ad> yolunu çözer; o dil yoksa diğerine düşer."""
    for d in (dil, "en" if dil == "tr" else "tr"):
        yol = os.path.join(GORSEL_KLASOR, d, ad)
        if os.path.exists(yol):
            return yol
    return None


def _gorsel_olcu(yol):
    """Piksel ölçüleri → (genişlik, yükseklik); okunamıyorsa 1600x920 varsay."""
    try:
        from PIL import Image
        with Image.open(yol) as im:
            return im.size
    except Exception:
        return (1600, 920)


def _gorsel_yerlestir(s, t, yol, x, y, w, h, dikey="orta"):
    """Görseli (x,y,w,h) kutusuna EN-BOY KORUYARAK yerleştirir, çerçeve çizer."""
    pw, ph = _gorsel_olcu(yol)
    oran = pw / ph
    if w / h > oran:                     # kutu görselden daha yayvan → yükseklik sınır
        gh, gw = h, Emu(int(h * oran))
    else:
        gw, gh = w, Emu(int(w / oran))
    gx = x + Emu(int((w - gw) / 2))
    gy = y if dikey == "ust" else y + Emu(int((h - gh) / 2))
    _dikdortgen(s, gx - Emu(19050), gy - Emu(19050), gw + Emu(38100),
                gh + Emu(38100), t.panel, t.cizgi, yaricap=False)
    s.shapes.add_picture(yol, gx, gy, gw, gh)
    return gx, gy, gw, gh


def _gorsel_eksik(s, t, ad, x, y, w, h):
    _dikdortgen(s, x, y, w, h, t.panel, t.cizgi)
    _metin(s, x, y + h / 2 - Inches(0.3), w, Inches(0.6),
           [(f"GÖRSEL EKSİK: {ad} — önce `python tanitim_ekran.py` çalıştırın",
             14, t.soluk, False)], hiza=PP_ALIGN.CENTER)


def _ciz_gorsel(s, t, veri, dil):
    """Tek büyük ekran görüntüsü: gerçek arayüz slaydın kendisidir."""
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    x, w = KENAR, GENISLIK - 2 * KENAR
    alt_pay = Inches(0.88) if veri.get("dipnot") else Inches(0.62)
    h = YUKSEKLIK - ust - alt_pay
    yol = _gorsel_yolu(veri["gorsel"], dil)
    if yol:
        _gorsel_yerlestir(s, t, yol, x, ust, w, h)
    else:
        _gorsel_eksik(s, t, veri["gorsel"], x, ust, w, h)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


def _ciz_galeri(s, t, veri, dil):
    """Yan yana iki ekran görüntüsü + altlarında birer başlık."""
    ust = _baslik_bandi(s, t, _c(veri["baslik"], dil), _c(veri.get("alt"), dil))
    ogeler = veri["ogeler"]
    n = len(ogeler)
    w = (GENISLIK - 2 * KENAR - Inches(0.35) * (n - 1)) / n
    h = YUKSEKLIK - ust - Inches(1.05)
    for i, (ad, altyazi) in enumerate(ogeler):
        x = KENAR + i * (w + Inches(0.35))
        yol = _gorsel_yolu(ad, dil)
        if yol:
            _, gy, _, gh = _gorsel_yerlestir(s, t, yol, x, ust + Inches(0.15),
                                             w, h, dikey="ust")
        else:
            _gorsel_eksik(s, t, ad, x, ust + Inches(0.15), w, h)
            gy, gh = ust + Inches(0.15), h
        _metin(s, x, gy + gh + Inches(0.18), w, Inches(0.35),
               [(_c(altyazi, dil), 13, t.soluk, False)], hiza=PP_ALIGN.CENTER)
    _dipnot(s, t, _c(veri.get("dipnot"), dil))


CIZICILER = {"kapak": _ciz_kapak, "madde": _ciz_madde, "kart": _ciz_kart,
             "sayi": _ciz_sayi, "tablo": _ciz_tablo, "kapanis": _ciz_kapanis,
             "gorsel": _ciz_gorsel, "galeri": _ciz_galeri}


def _c(cift, dil):
    """(tr, en) çiftinden dili seçer; tek dize verilmişse olduğu gibi döner."""
    if cift is None:
        return None
    if isinstance(cift, (tuple, list)):
        return cift[1] if (dil == "en" and len(cift) > 1 and cift[1]) else cift[0]
    return cift


# ══════════════════════════════════════════════════════════════════════════
# ÜRETİM
# ══════════════════════════════════════════════════════════════════════════
def sunum_kur(dil, mod):
    """Sıfırdan bir Presentation kurar → (prs, slayt_sayisi)."""
    t = _Tasarim(mod)
    prs = Presentation()
    prs.slide_width = GENISLIK
    prs.slide_height = YUKSEKLIK
    toplam = len(SLAYTLAR)
    for i, veri in enumerate(SLAYTLAR, 1):
        s = _slayt(prs, t)
        CIZICILER[veri["tip"]](s, t, veri, dil)
        if veri["tip"] not in ("kapak", "kapanis"):
            _sayfa_no(s, t, i, toplam)
    return prs, toplam


def _powerpoint():
    try:
        import win32com.client
        return win32com.client.Dispatch("PowerPoint.Application")
    except Exception:
        return None


def pdf_ve_png_uret(pptx_yolu, pdf_yolu, png_klasoru=None):
    """PowerPoint COM ile PDF (+ istenirse slayt PNG'leri) üretir → True/False."""
    app = _powerpoint()
    if app is None:
        return False
    sunum = None
    try:
        sunum = app.Presentations.Open(pptx_yolu, WithWindow=False)
        sunum.SaveAs(pdf_yolu, 32)                       # 32 = ppSaveAsPDF
        if png_klasoru:
            os.makedirs(png_klasoru, exist_ok=True)
            sunum.SaveAs(png_klasoru, 18)                # 18 = ppSaveAsPNG
        return True
    except Exception as exc:
        print(f"   [!] PowerPoint dışa aktarma hatası: {exc}")
        return False
    finally:
        try:
            if sunum is not None:
                sunum.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass


HTML_SABLON = """<!DOCTYPE html>
<html lang="__LANG__"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>__TITLE__</title>
<style>
  :root{--bg:#201B15;--panel:#2A241C;--ink:#EDE5D8;--mut:#A6987F;--amber:#E07B39;}
  *{box-sizing:border-box;}
  body{margin:0;background:var(--bg);color:var(--ink);
    font-family:'__FONT__',system-ui,-apple-system,Arial,sans-serif;}
  header{padding:18px 22px;display:flex;align-items:center;gap:12px;
    border-bottom:1px solid #3A3226;position:sticky;top:0;background:var(--bg);z-index:5;}
  header .ic{width:34px;height:34px;border-radius:9px;background:var(--amber);
    display:flex;align-items:center;justify-content:center;font-size:19px;font-weight:700;color:#fff;}
  header b{font-size:15px;letter-spacing:.4px;}
  header .sub{color:var(--mut);font-size:12.5px;margin-left:auto;}
  .stage{max-width:1180px;margin:0 auto;padding:18px;}
  .slide{margin:0 0 22px;border:1px solid #3A3226;border-radius:12px;overflow:hidden;
    background:var(--panel);}
  .slide img{display:block;width:100%;height:auto;}
  .slide .no{padding:8px 12px;color:var(--mut);font-size:12px;
    border-top:1px solid #3A3226;}
  .txt{padding:20px 22px;}
  .txt h2{margin:0 0 6px;font-size:21px;color:var(--ink);}
  .txt p{margin:6px 0;color:var(--mut);font-size:14px;line-height:1.5;}
  .txt li{color:var(--ink);font-size:14px;line-height:1.55;margin:5px 0;}
  nav{position:fixed;right:16px;bottom:16px;display:flex;gap:8px;z-index:9;}
  nav button{font:inherit;font-size:16px;font-weight:700;width:44px;height:44px;
    border-radius:11px;border:1px solid #3A3226;background:var(--panel);
    color:var(--ink);cursor:pointer;}
  nav button:hover{background:var(--amber);color:#fff;}
  footer{color:var(--mut);font-size:12.5px;padding:10px 22px 40px;text-align:center;}
  @media print{nav{display:none;} body{background:#fff;color:#000;}
    .slide{break-inside:avoid;border-color:#ccc;background:#fff;}}
</style></head>
<body>
<header><span class="ic">Σ</span><b>MATH COURSE AI</b><span class="sub">__SUB__</span></header>
<div class="stage">__SLIDES__</div>
<nav><button onclick="git(-1)">↑</button><button onclick="git(1)">↓</button></nav>
<footer>__FOOT__</footer>
<script>
  var i = 0;
  function git(d){
    var s = document.querySelectorAll('.slide');
    i = Math.max(0, Math.min(s.length - 1, i + d));
    s[i].scrollIntoView({behavior:'smooth', block:'start'});
  }
  document.addEventListener('keydown', function(e){
    if (e.key === 'ArrowDown' || e.key === 'PageDown' || e.key === ' ') { git(1); e.preventDefault(); }
    if (e.key === 'ArrowUp' || e.key === 'PageUp') { git(-1); e.preventDefault(); }
  });
</script>
</body></html>"""


def _slayt_metni_html(veri, dil):
    """PNG yoksa slaydı METİN olarak yazar — HTML her hâlükârda okunur olsun."""
    e = _html.escape
    par = [f'<h2>{e(_c(veri["baslik"], dil) or "")}</h2>']
    alt = _c(veri.get("alt"), dil)
    if alt:
        par.append(f"<p>{e(alt)}</p>")
    tip = veri["tip"]
    ogeler = veri.get("ogeler") or []
    if tip in ("madde", "kapanis"):
        par.append("<ul>" + "".join(f"<li>{e(_c(o, dil))}</li>" for o in ogeler) + "</ul>")
    elif tip == "kart":
        par.append("<ul>" + "".join(
            f"<li><b>{e(_c(b, dil))}</b> — {e(_c(g, dil))}</li>" for b, g in ogeler) + "</ul>")
    elif tip == "sayi":
        par.append("<ul>" + "".join(
            f"<li><b>{e(_c(s, dil))}</b> — {e(_c(a, dil).replace(chr(10), ' '))}</li>"
            for s, a in ogeler) + "</ul>")
    elif tip == "tablo":
        par.append("<ul>" + "".join(
            f"<li><b>{e(_c(s, dil))}</b> → {e(_c(g, dil))}</li>" for s, g in ogeler) + "</ul>")
    elif tip == "gorsel":
        par.append(f"<p>[{e(veri['gorsel'])} — ekran görüntüsü]</p>")
    elif tip == "galeri":
        par.append("<ul>" + "".join(
            f"<li>[{e(ad)}] {e(_c(altyazi, dil))}</li>" for ad, altyazi in ogeler) + "</ul>")
    dip = _c(veri.get("dipnot"), dil)
    if dip:
        par.append(f"<p>{e(dip)}</p>")
    return '<div class="txt">' + "".join(par) + "</div>"


def html_uret(yol, dil, png_yollari):
    """Slayt PNG'lerini GÖMEN tek dosyalık sunum; PNG yoksa metne düşer."""
    parcalar = []
    for i, veri in enumerate(SLAYTLAR):
        ic = ""
        if i < len(png_yollari):
            try:
                with open(png_yollari[i], "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("ascii")
                ic = f'<img src="data:image/png;base64,{b64}" alt="{i + 1}">'
            except Exception:
                ic = ""
        if not ic:
            ic = _slayt_metni_html(veri, dil)
        parcalar.append(f'<section class="slide">{ic}'
                        f'<div class="no">{i + 1} / {len(SLAYTLAR)}</div></section>')
    baslik = ("Math Course AI — Tanıtım" if dil == "tr" else "Math Course AI — Intro")
    foot = (f"{VERSION} · bu sayfa programla birlikte üretildi, çevrimdışı çalışır"
            if dil == "tr" else
            f"{VERSION} · generated with the program, works offline")
    metin = (HTML_SABLON.replace("__FONT__", YAZI_TIPI)
             .replace("__LANG__", dil).replace("__TITLE__", baslik)
             .replace("__SUB__", _c(ALTBASLIK, dil)).replace("__FOOT__", foot)
             .replace("__SLIDES__", "".join(parcalar)))
    with open(yol, "w", encoding="utf-8") as f:
        f.write(metin)
    return len(metin)


def _png_sirala(klasor):
    dosyalar = glob.glob(os.path.join(klasor, "*.PNG")) + \
        glob.glob(os.path.join(klasor, "*.png"))

    def anahtar(p):
        ad = os.path.splitext(os.path.basename(p))[0]
        rakam = "".join(ch for ch in ad if ch.isdigit())
        return int(rakam) if rakam else 0
    return sorted(set(dosyalar), key=anahtar)


def uret(sadece_pptx=False):
    sonuc = []
    png_onbellek = {}
    for (dil, mod), ad in CIKTILAR.items():
        pptx_yolu = os.path.join(KOK, ad + ".pptx")
        prs, n = sunum_kur(dil, mod)
        prs.save(pptx_yolu)
        satir = {"ad": ad, "dil": dil, "mod": mod, "pptx": True, "slayt": n,
                 "pdf": False, "html": None}
        if not sadece_pptx:
            pdf_yolu = os.path.join(KOK, ad + ".pdf")
            png_klasor = None
            if mod == "ekran":
                png_klasor = os.path.join(tempfile.mkdtemp(prefix="mca_slayt_"), "png")
            satir["pdf"] = pdf_ve_png_uret(pptx_yolu, pdf_yolu, png_klasor)
            if mod == "ekran":
                pngler = _png_sirala(png_klasor) if png_klasor else []
                png_onbellek[dil] = pngler
                html_yolu = os.path.join(KOK, ad + ".html")
                satir["html"] = html_uret(html_yolu, dil, pngler)
                satir["png"] = len(pngler)
                if png_klasor:
                    shutil.rmtree(os.path.dirname(png_klasor), ignore_errors=True)
        sonuc.append(satir)
    return sonuc


def kontrol():
    print(f"{'dosya':44} {'durum':>10} {'boyut':>12}")
    print("-" * 70)
    for (dil, mod), ad in CIKTILAR.items():
        for uzanti in (".pptx", ".pdf", ".html"):
            if mod == "baski" and uzanti == ".html":
                continue
            yol = os.path.join(KOK, ad + uzanti)
            if os.path.exists(yol):
                print(f"{ad + uzanti:44} {'VAR':>10} {os.path.getsize(yol):>12,}")
            else:
                print(f"{ad + uzanti:44} {'YOK':>10} {'-':>12}")


def main():
    if "--kontrol" in sys.argv:
        kontrol()
        return 0
    sadece = "--sadece-pptx" in sys.argv
    print(f"Math Course AI — tanıtım sunumu üretiliyor ({VERSION})")
    print(f"{len(SLAYTLAR)} slayt · {len(CIKTILAR)} çıktı seti"
          + ("  [yalnız PPTX]" if sadece else ""))
    print("-" * 70)
    satirlar = uret(sadece_pptx=sadece)
    eksik = 0
    for s in satirlar:
        parcalar = [f"pptx({s['slayt']} slayt)"]
        if not sadece:
            parcalar.append("pdf" if s["pdf"] else "PDF YOK")
            if s["html"] is not None:
                parcalar.append(f"html({s['html']:,} bayt, {s.get('png', 0)} png)")
        if not sadece and not s["pdf"]:
            eksik += 1
        print(f"  {s['ad']:38} {' · '.join(parcalar)}")
    print("-" * 70)
    if eksik:
        print(f"[!] {eksik} PDF üretilemedi — PowerPoint kurulu mu? PPTX ve HTML hazır.")
        return 1
    print("[OK] Hepsi üretildi.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
